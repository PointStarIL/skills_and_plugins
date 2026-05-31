# -*- coding: utf-8 -*-
"""HebrewLegalDeck — גרסה מלאה ומורחבת עם שיפורים גרפיים."""

from __future__ import annotations
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from lxml import etree

# ====================================================================
# Design System
# ====================================================================
NAVY = RGBColor(0x0E, 0x2A, 0x47)
IVORY = RGBColor(0xF5, 0xF1, 0xEA)
IVORY_DARK = RGBColor(0xEB, 0xE5, 0xDB)
NEAR_BLACK = RGBColor(0x1F, 0x1F, 0x1F)
GOLD = RGBColor(0xB8, 0x86, 0x0B)
GOLD_LIGHT = RGBColor(0xD4, 0xA5, 0x2E)
SLATE = RGBColor(0x5A, 0x6B, 0x7D)

TITLE_FONT = "Suez One"
BODY_FONT = "Rubik"

TITLE_PT = 36
BODY_PT = 24
FOOTER_PT = 12

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_X = Inches(0.6)
MARGIN_TOP = Inches(0.55)
MARGIN_BOTTOM = Inches(0.4)
TOP_LINE_HEIGHT = Pt(3)
GOLD_SQUARE_SIZE = Inches(0.18)
GOLD_SQUARE_MARGIN = Inches(0.25)

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _qn(tag):
    return f"{{{A_NS}}}{tag}"


def set_paragraph_rtl(paragraph):
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")
    paragraph.alignment = PP_ALIGN.RIGHT


def set_run_lang_he(run):
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "he-IL")


def set_run_complex_font(run, typeface):
    rPr = run._r.get_or_add_rPr()
    cs_tag = f"{{{A_NS}}}cs"
    for existing in rPr.findall(cs_tag):
        rPr.remove(existing)
    cs = etree.SubElement(rPr, cs_tag)
    cs.set("typeface", typeface)
    ea_tag = f"{{{A_NS}}}ea"
    for existing in rPr.findall(ea_tag):
        rPr.remove(existing)
    ea = etree.SubElement(rPr, ea_tag)
    ea.set("typeface", typeface)


def format_numbered_text(number, text):
    return f"{number}. {text}"


# ====================================================================
# HebrewLegalDeck
# ====================================================================
class HebrewLegalDeck:

    def __init__(self, lecturer_name, output_path,
                 title_font=TITLE_FONT, body_font=BODY_FONT):
        self.lecturer_name = lecturer_name
        self.output_path = Path(output_path)
        self.title_font = title_font
        self.body_font = body_font
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self._slide_count = 0

    # ------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------
    def add_title_slide(self, title, subtitle=None, image_path=None):
        slide = self._new_slide(include_footer=False)

        # אם יש תמונה — נציג ברקע עם שכבת shading
        if image_path:
            try:
                from PIL import Image
                with Image.open(str(image_path)) as im:
                    src_w, src_h = im.size
            except Exception:
                src_w, src_h = 16, 9

            # תמונה בצד שמאל של השקף, לוקחת חצי משקף
            img_area_w = SLIDE_WIDTH // 2
            img_area_h = SLIDE_HEIGHT
            src_ratio = src_w / src_h
            avail_ratio = img_area_w / img_area_h

            if src_ratio > avail_ratio:
                final_w = img_area_w
                final_h = int(img_area_w / src_ratio)
            else:
                final_h = img_area_h
                final_w = int(img_area_h * src_ratio)

            img_left = (img_area_w - final_w) // 2
            img_top = (img_area_h - final_h) // 2
            slide.shapes.add_picture(
                str(image_path),
                img_left, img_top,
                width=final_w, height=final_h,
            )

            # רצועה זהובה אנכית מפרידה בין התמונה לטקסט
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                img_area_w - Inches(0.04), 0,
                Inches(0.08), SLIDE_HEIGHT,
            )
            divider.line.fill.background()
            divider.fill.solid()
            divider.fill.fore_color.rgb = GOLD

            # טקסט בצד ימין
            text_x = img_area_w + Inches(0.4)
            text_w = SLIDE_WIDTH - text_x - Inches(0.6)
        else:
            text_x = MARGIN_X
            text_w = SLIDE_WIDTH - 2 * MARGIN_X

        # כותרת מרכזית
        title_box = slide.shapes.add_textbox(
            text_x, Inches(2.4), text_w, Inches(2.0),
        )
        self._fill_text(
            title_box, title,
            font_name=self.title_font, font_size=42 if image_path else 48,
            color=NAVY, bold=True,
        )

        if subtitle:
            sub_box = slide.shapes.add_textbox(
                text_x, Inches(4.5), text_w, Inches(0.8),
            )
            self._fill_text(
                sub_box, subtitle,
                font_name=self.body_font, font_size=22,
                color=SLATE,
            )

        # שם המרצה
        lec_box = slide.shapes.add_textbox(
            text_x, SLIDE_HEIGHT - Inches(0.9),
            text_w, Inches(0.5),
        )
        self._fill_text(
            lec_box, self.lecturer_name,
            font_name=self.body_font, font_size=14,
            color=SLATE,
        )

    def add_content_slide(self, title, bullets):
        slide = self._new_slide()
        self._add_title(slide, title)
        self._add_decorative_underline(slide)
        self._add_bullets(slide, bullets, numbered=False)

    def add_numbered_slide(self, title, items):
        slide = self._new_slide()
        self._add_title(slide, title)
        self._add_decorative_underline(slide)
        self._add_bullets(slide, items, numbered=True)

    def add_two_column_slide(self, title, left_title, left_items,
                             right_title, right_items):
        slide = self._new_slide()
        self._add_title(slide, title)
        self._add_decorative_underline(slide)

        column_width = (SLIDE_WIDTH - 2 * MARGIN_X - Inches(0.4)) // 2
        col_top = Inches(1.7)
        col_height = SLIDE_HEIGHT - col_top - Inches(0.8)

        right_x = SLIDE_WIDTH - MARGIN_X - column_width
        self._add_column(slide, right_x, col_top, column_width, col_height,
                         right_title, right_items)
        left_x = MARGIN_X
        self._add_column(slide, left_x, col_top, column_width, col_height,
                         left_title, left_items)

    def add_quote_slide(self, quote, source):
        slide = self._new_slide()
        # סמל ציטוט גדול בזהב
        quote_mark_box = slide.shapes.add_textbox(
            SLIDE_WIDTH - Inches(2), Inches(1.0),
            Inches(1.5), Inches(1.5),
        )
        tf = quote_mark_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = '"'
        run.font.name = self.title_font
        run.font.size = Pt(120)
        run.font.color.rgb = GOLD
        run.font.bold = True

        quote_box = slide.shapes.add_textbox(
            MARGIN_X + Inches(0.5), Inches(2.6),
            SLIDE_WIDTH - 2 * (MARGIN_X + Inches(0.5)), Inches(2.8),
        )
        self._fill_text(
            quote_box, quote,
            font_name=self.title_font, font_size=24,
            color=NEAR_BLACK, italic=True,
        )

        source_box = slide.shapes.add_textbox(
            MARGIN_X + Inches(0.5), Inches(5.6),
            SLIDE_WIDTH - 2 * (MARGIN_X + Inches(0.5)), Inches(0.5),
        )
        self._fill_text(
            source_box, f"— {source}",
            font_name=self.body_font, font_size=16,
            color=GOLD,
        )

    def add_section_divider(self, title, chapter_num=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._slide_count += 1

        # רקע נייבי
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT,
        )
        bg.line.fill.background()
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

        # רצועה זהובה אופקית עליונה
        gold_band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(2.5),
            Inches(2.5), Inches(0.08),
        )
        gold_band.line.fill.background()
        gold_band.fill.solid()
        gold_band.fill.fore_color.rgb = GOLD

        # ריבוע זהב גדול
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            SLIDE_WIDTH - Inches(2.0), Inches(3.0),
            Inches(1.5), Inches(1.5),
        )
        accent.line.fill.background()
        accent.fill.solid()
        accent.fill.fore_color.rgb = GOLD

        # קו דק זהוב מתחת לכותרת
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN_X, Inches(5.2),
            Inches(2.0), Pt(2),
        )
        title_underline.line.fill.background()
        title_underline.fill.solid()
        title_underline.fill.fore_color.rgb = GOLD

        # טקסט כותרת
        title_box = slide.shapes.add_textbox(
            MARGIN_X, Inches(2.7),
            SLIDE_WIDTH - Inches(2.5) - MARGIN_X, Inches(2.5),
        )
        self._fill_text(
            title_box, title,
            font_name=self.title_font, font_size=46,
            color=IVORY, bold=True,
        )

        # מספר שקף
        page_box = slide.shapes.add_textbox(
            SLIDE_WIDTH - Inches(0.8), SLIDE_HEIGHT - Inches(0.5),
            Inches(0.5), Inches(0.3),
        )
        self._fill_text(
            page_box, str(self._slide_count),
            font_name=self.body_font, font_size=18,
            color=GOLD,
        )

    def add_blank_slide(self):
        return self._new_slide()

    def add_image_full_slide(self, image_path, title=None, caption=None):
        slide = self._new_slide()
        from pptx.util import Inches as _Inches

        if title:
            self._add_title(slide, title)
            self._add_decorative_underline(slide)
            img_top = _Inches(1.7)
            img_height_max = SLIDE_HEIGHT - img_top - _Inches(1.2)
        else:
            img_top = _Inches(0.7)
            img_height_max = SLIDE_HEIGHT - img_top - _Inches(0.9)

        if caption:
            img_height_max -= _Inches(0.5)

        from PIL import Image
        try:
            with Image.open(str(image_path)) as im:
                src_w, src_h = im.size
        except Exception:
            src_w, src_h = 16, 9

        avail_w = SLIDE_WIDTH - 2 * MARGIN_X
        avail_h = img_height_max

        src_ratio = src_w / src_h
        avail_ratio = avail_w / avail_h

        if src_ratio > avail_ratio:
            final_w = avail_w
            final_h = int(avail_w / src_ratio)
        else:
            final_h = avail_h
            final_w = int(avail_h * src_ratio)

        img_left = (SLIDE_WIDTH - final_w) // 2
        slide.shapes.add_picture(
            str(image_path), img_left, img_top,
            width=final_w, height=final_h,
        )

        if caption:
            cap_box = slide.shapes.add_textbox(
                MARGIN_X, img_top + final_h + _Inches(0.15),
                SLIDE_WIDTH - 2 * MARGIN_X, _Inches(0.4),
            )
            self._fill_text(
                cap_box, caption,
                font_name=self.body_font, font_size=14,
                color=SLATE,
            )

    def add_content_slide_with_image(self, title, bullets, image_path,
                                     image_side="left"):
        from pptx.util import Inches as _Inches
        slide = self._new_slide()
        self._add_title(slide, title)
        self._add_decorative_underline(slide)

        body_top = _Inches(1.75)
        body_height = SLIDE_HEIGHT - body_top - _Inches(0.6)
        gap = _Inches(0.3)
        col_width = (SLIDE_WIDTH - 2 * MARGIN_X - gap) // 2

        if image_side == "left":
            img_x = MARGIN_X
            text_x = MARGIN_X + col_width + gap
        else:
            text_x = MARGIN_X
            img_x = MARGIN_X + col_width + gap

        # תיבת טקסט
        text_box = slide.shapes.add_textbox(
            text_x, body_top, col_width, body_height,
        )
        tf = text_box.text_frame
        tf.word_wrap = True

        n = len(bullets)
        avg_len = sum(len(s) for s in bullets) / max(n, 1)
        if avg_len > 100:
            body_pt = 14
        elif avg_len > 70:
            body_pt = 16
        else:
            body_pt = 18

        for i, item in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = f"•  {item}"
            run.font.name = self.body_font
            run.font.size = Pt(body_pt)
            run.font.color.rgb = NEAR_BLACK
            set_run_lang_he(run)
            set_run_complex_font(run, self.body_font)
            set_paragraph_rtl(p)
            p.alignment = PP_ALIGN.RIGHT

        # תמונה בצד השני
        from PIL import Image
        try:
            with Image.open(str(image_path)) as im:
                src_w, src_h = im.size
        except Exception:
            src_w, src_h = 1, 1

        src_ratio = src_w / src_h
        avail_ratio = col_width / body_height

        if src_ratio > avail_ratio:
            final_w = col_width
            final_h = int(col_width / src_ratio)
        else:
            final_h = body_height
            final_w = int(body_height * src_ratio)

        img_y = body_top + (body_height - final_h) // 2
        img_x_centered = img_x + (col_width - final_w) // 2

        slide.shapes.add_picture(
            str(image_path), img_x_centered, img_y,
            width=final_w, height=final_h,
        )

    def save(self, path=None):
        target = Path(path) if path else self.output_path
        target.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(target))
        return target

    # ------------------------------------------------------------
    # Internal building blocks
    # ------------------------------------------------------------
    def _new_slide(self, include_footer=True):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._slide_count += 1

        # רקע שנהב
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT,
        )
        bg.line.fill.background()
        bg.fill.solid()
        bg.fill.fore_color.rgb = IVORY
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

        # קו נייבי בראש
        top_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, TOP_LINE_HEIGHT,
        )
        top_line.line.fill.background()
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = NAVY

        # רצועה דקה זהובה מתחת לקו הנייבי
        gold_strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, TOP_LINE_HEIGHT, Inches(2.5), Pt(1.5),
        )
        gold_strip.line.fill.background()
        gold_strip.fill.solid()
        gold_strip.fill.fore_color.rgb = GOLD

        # ריבוע זהב בפינה
        gold_x = SLIDE_WIDTH - GOLD_SQUARE_MARGIN - GOLD_SQUARE_SIZE
        gold_y = GOLD_SQUARE_MARGIN
        gold_sq = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            gold_x, gold_y, GOLD_SQUARE_SIZE, GOLD_SQUARE_SIZE,
        )
        gold_sq.line.fill.background()
        gold_sq.fill.solid()
        gold_sq.fill.fore_color.rgb = GOLD

        # קו פוטר דק
        footer_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN_X, SLIDE_HEIGHT - Inches(0.55),
            SLIDE_WIDTH - 2 * MARGIN_X, Pt(0.5),
        )
        footer_line.line.fill.background()
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = SLATE

        if include_footer:
            self._add_footer(slide)
        return slide

    def _add_decorative_underline(self, slide):
        """קו זהב דק מתחת לכותרת השקף — שיפור עיצובי."""
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            SLIDE_WIDTH - MARGIN_X - Inches(2.5),
            Inches(1.45),
            Inches(2.5), Pt(2.5),
        )
        underline.line.fill.background()
        underline.fill.solid()
        underline.fill.fore_color.rgb = GOLD

    def _add_footer(self, slide):
        footer_y = SLIDE_HEIGHT - Inches(0.4)
        name_box = slide.shapes.add_textbox(
            MARGIN_X, footer_y, Inches(8), Inches(0.3),
        )
        tf = name_box.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = self.lecturer_name
        run.font.name = self.body_font
        run.font.size = Pt(FOOTER_PT)
        run.font.color.rgb = SLATE
        set_run_lang_he(run)
        set_run_complex_font(run, self.body_font)
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        pPr.set("rtl", "1")
        p.alignment = PP_ALIGN.LEFT

        num_box = slide.shapes.add_textbox(
            SLIDE_WIDTH - MARGIN_X - Inches(1),
            footer_y, Inches(1), Inches(0.3),
        )
        tf2 = num_box.text_frame
        tf2.margin_left = 0
        tf2.margin_right = 0
        tf2.margin_top = 0
        tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        run2 = p2.add_run()
        run2.text = str(self._slide_count)
        run2.font.name = self.body_font
        run2.font.size = Pt(FOOTER_PT)
        run2.font.color.rgb = SLATE
        set_run_complex_font(run2, self.body_font)

    def _add_title(self, slide, title):
        title_box = slide.shapes.add_textbox(
            MARGIN_X, MARGIN_TOP,
            SLIDE_WIDTH - 2 * MARGIN_X, Inches(0.9),
        )
        self._fill_text(
            title_box, title,
            font_name=self.title_font, font_size=TITLE_PT,
            color=NAVY, bold=True,
        )

    def _add_bullets(self, slide, items, numbered=False):
        body_top = Inches(1.7)
        body_height = SLIDE_HEIGHT - body_top - Inches(0.6)
        body_box = slide.shapes.add_textbox(
            MARGIN_X, body_top,
            SLIDE_WIDTH - 2 * MARGIN_X, body_height,
        )
        tf = body_box.text_frame
        tf.word_wrap = True

        n = len(items)
        avg_len = sum(len(s) for s in items) / max(n, 1)
        if n >= 5 and avg_len > 80:
            body_pt = 18
        elif n >= 5:
            body_pt = 20
        else:
            body_pt = 22

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(10)
            text = format_numbered_text(i + 1, item) if numbered else f"•  {item}"
            run = p.add_run()
            run.text = text
            run.font.name = self.body_font
            run.font.size = Pt(body_pt)
            run.font.color.rgb = NEAR_BLACK
            set_run_lang_he(run)
            set_run_complex_font(run, self.body_font)
            set_paragraph_rtl(p)

    def _add_column(self, slide, x, y, width, height, col_title, items):
        col_title_box = slide.shapes.add_textbox(x, y, width, Inches(0.5))
        self._fill_text(
            col_title_box, col_title,
            font_name=self.title_font, font_size=20,
            color=GOLD, bold=True,
        )
        items_box = slide.shapes.add_textbox(
            x, y + Inches(0.6), width, height - Inches(0.6),
        )
        tf = items_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"•  {item}"
            run.font.name = self.body_font
            run.font.size = Pt(BODY_PT)
            run.font.color.rgb = NEAR_BLACK
            set_run_lang_he(run)
            set_run_complex_font(run, self.body_font)
            set_paragraph_rtl(p)

    def _fill_text(self, text_box, text, font_name, font_size, color,
                   bold=False, italic=False):
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        set_run_lang_he(run)
        set_run_complex_font(run, font_name)
        set_paragraph_rtl(p)
