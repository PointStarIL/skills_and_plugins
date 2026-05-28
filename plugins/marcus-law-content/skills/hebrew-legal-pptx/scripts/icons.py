# -*- coding: utf-8 -*-
"""ספריית איקונים גיאומטריים מצוירים ישירות עם python-pptx shapes."""

from pptx.util import Pt, Inches, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

NAVY = RGBColor(0x0E, 0x2A, 0x47)
GOLD = RGBColor(0xB8, 0x86, 0x0B)
SLATE = RGBColor(0x5A, 0x6B, 0x7D)
IVORY = RGBColor(0xF5, 0xF1, 0xEA)


def _add_shape(slide, shape_type, x, y, w, h, fill_color=None,
               line_color=None, line_width=None):
    """Helper להוספת shape עם הגדרות בסיסיות."""
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        if line_width:
            shp.line.width = line_width
    if fill_color is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    return shp


def icon_gavel(slide, cx, cy, size=Inches(1.2), color=NAVY, accent=GOLD):
    """איקון פטיש שופט — פטיש מוטה אלכסונית עם בסיס."""
    # ראש הפטיש - מלבן עבה בזווית
    head_w = size
    head_h = int(size * 0.35)
    head_x = cx - head_w // 2
    head_y = cy - int(size * 0.6)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, head_x, head_y, head_w, head_h,
               fill_color=color)
    # ידית
    handle_w = int(size * 0.6)
    handle_h = int(size * 0.18)
    handle_x = cx - handle_w // 2
    handle_y = head_y + head_h + int(size * 0.05)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, handle_x, handle_y, handle_w, handle_h,
               fill_color=color)
    # בסיס מתחת
    base_w = int(size * 1.1)
    base_h = int(size * 0.12)
    base_x = cx - base_w // 2
    base_y = handle_y + handle_h + int(size * 0.1)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, base_x, base_y, base_w, base_h,
               fill_color=accent)


def icon_scales(slide, cx, cy, size=Inches(1.2), color=NAVY, accent=GOLD):
    """איקון מאזני צדק — קו אופקי עם שני צלחות."""
    # עמוד אנכי
    pole_w = int(size * 0.06)
    pole_h = size
    pole_x = cx - pole_w // 2
    pole_y = cy - size // 2
    _add_shape(slide, MSO_SHAPE.RECTANGLE, pole_x, pole_y, pole_w, pole_h,
               fill_color=color)
    # קורה אופקית
    beam_w = int(size * 1.2)
    beam_h = int(size * 0.06)
    beam_x = cx - beam_w // 2
    beam_y = cy - size // 2 + int(size * 0.1)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, beam_x, beam_y, beam_w, beam_h,
               fill_color=color)
    # צלחת ימין
    plate_w = int(size * 0.45)
    plate_h = int(size * 0.12)
    _add_shape(slide, MSO_SHAPE.OVAL,
               beam_x + beam_w - plate_w, beam_y + beam_h + int(size * 0.18),
               plate_w, plate_h, fill_color=accent)
    # צלחת שמאל
    _add_shape(slide, MSO_SHAPE.OVAL,
               beam_x, beam_y + beam_h + int(size * 0.18),
               plate_w, plate_h, fill_color=accent)
    # בסיס
    base_w = int(size * 0.5)
    base_h = int(size * 0.08)
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               cx - base_w // 2, pole_y + pole_h - base_h,
               base_w, base_h, fill_color=color)


def icon_document(slide, cx, cy, size=Inches(1.2), color=NAVY, accent=GOLD):
    """איקון מסמך — מלבן אנכי עם פינה מקופלת וקווי טקסט."""
    doc_w = int(size * 0.75)
    doc_h = size
    doc_x = cx - doc_w // 2
    doc_y = cy - doc_h // 2
    # גוף המסמך
    _add_shape(slide, MSO_SHAPE.RECTANGLE, doc_x, doc_y, doc_w, doc_h,
               fill_color=IVORY, line_color=color, line_width=Pt(2))
    # פינה מקופלת
    fold_size = int(size * 0.18)
    _add_shape(slide, MSO_SHAPE.RIGHT_TRIANGLE,
               doc_x + doc_w - fold_size, doc_y, fold_size, fold_size,
               fill_color=color)
    # שלושה קווי טקסט
    line_w = int(doc_w * 0.7)
    line_h = int(size * 0.04)
    line_x = doc_x + (doc_w - line_w) // 2
    for i in range(3):
        line_y = doc_y + int(size * 0.35) + i * int(size * 0.18)
        _add_shape(slide, MSO_SHAPE.RECTANGLE, line_x, line_y, line_w, line_h,
                   fill_color=color if i < 2 else accent)


def icon_building(slide, cx, cy, size=Inches(1.2), color=NAVY, accent=GOLD):
    """איקון מוסד תכנון — בנין עם עמודים."""
    bld_w = size
    bld_h = int(size * 0.85)
    bld_x = cx - bld_w // 2
    bld_y = cy - bld_h // 2
    # גג משולש
    triangle_h = int(size * 0.2)
    _add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE,
               bld_x, bld_y, bld_w, triangle_h, fill_color=accent)
    # קורה עליונה
    beam_h = int(size * 0.06)
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               bld_x, bld_y + triangle_h, bld_w, beam_h, fill_color=color)
    # 4 עמודים
    pillars_y = bld_y + triangle_h + beam_h
    pillars_h = int(size * 0.45)
    pillar_w = int(bld_w / 9)
    gap = (bld_w - 4 * pillar_w) // 5
    for i in range(4):
        x = bld_x + gap + i * (pillar_w + gap)
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   x, pillars_y, pillar_w, pillars_h, fill_color=color)
    # בסיס
    base_y = pillars_y + pillars_h
    base_h = int(size * 0.07)
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               bld_x, base_y, bld_w, base_h, fill_color=color)


def icon_people(slide, cx, cy, size=Inches(1.2), color=NAVY, accent=GOLD):
    """איקון 3 דמויות עומדות יחד."""
    # 3 דמויות
    head_size = int(size * 0.18)
    body_w = int(size * 0.22)
    body_h = int(size * 0.45)
    spacing = int(size * 0.32)

    positions = [
        (cx - spacing, cy, color),
        (cx, cy - int(size * 0.05), accent),
        (cx + spacing, cy, color),
    ]
    for px, py, c in positions:
        # ראש
        _add_shape(slide, MSO_SHAPE.OVAL,
                   px - head_size // 2, py - body_h // 2 - head_size,
                   head_size, head_size, fill_color=c)
        # גוף
        _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                   px - body_w // 2, py - body_h // 2,
                   body_w, body_h, fill_color=c)


def icon_flow(slide, cx, cy, size=Inches(2.5), color=NAVY, accent=GOLD):
    """איקון תהליך זרימה - ארבעה צמתים מחוברים בקו אנכי."""
    node_size = int(size * 0.18)
    flow_h = size
    line_x = cx
    line_top = cy - flow_h // 2
    # קו אנכי
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               line_x - Inches(0.02), line_top, Inches(0.04), flow_h,
               fill_color=SLATE)
    # 4 צמתים
    n = 4
    for i in range(n):
        ny = line_top + (flow_h * i) // (n - 1) - node_size // 2
        c = accent if i in (0, n-1) else color
        _add_shape(slide, MSO_SHAPE.OVAL,
                   line_x - node_size // 2, ny, node_size, node_size,
                   fill_color=c)


def icon_map_pin(slide, cx, cy, size=Inches(1.2), color=NAVY, accent=GOLD):
    """איקון מפה עם פין מיקום."""
    # רקע מפה - מלבן מעוגל
    map_w = int(size * 1.1)
    map_h = int(size * 0.85)
    map_x = cx - map_w // 2
    map_y = cy - map_h // 2
    _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, map_x, map_y, map_w, map_h,
               fill_color=IVORY, line_color=SLATE, line_width=Pt(1.5))
    # קווי גריד דקיקים
    for i in range(1, 3):
        gx = map_x + (map_w * i) // 3
        _add_shape(slide, MSO_SHAPE.RECTANGLE, gx, map_y, Pt(0.5), map_h,
                   fill_color=SLATE)
        gy = map_y + (map_h * i) // 3
        _add_shape(slide, MSO_SHAPE.RECTANGLE, map_x, gy, map_w, Pt(0.5),
                   fill_color=SLATE)
    # פין מיקום במרכז
    pin_size = int(size * 0.32)
    _add_shape(slide, MSO_SHAPE.OVAL,
               cx - pin_size // 2, cy - pin_size // 2,
               pin_size, pin_size, fill_color=accent)
    _add_shape(slide, MSO_SHAPE.OVAL,
               cx - pin_size // 4, cy - pin_size // 4,
               pin_size // 2, pin_size // 2, fill_color=color)


def chapter_badge(slide, cx, cy, number, size=Inches(1.4),
                  color=GOLD, text_color=IVORY):
    """תג מספר פרק - עיגול גדול עם המספר במרכז."""
    from pptx.enum.text import PP_ALIGN
    from pptx_builder_complete import set_run_complex_font, set_paragraph_rtl
    # עיגול חיצוני
    _add_shape(slide, MSO_SHAPE.OVAL,
               cx - size // 2, cy - size // 2, size, size,
               fill_color=color)
    # טקסט המספר
    txt_box = slide.shapes.add_textbox(
        cx - size // 2, cy - size // 2,
        size, size,
    )
    tf = txt_box.text_frame
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.name = "Suez One"
    run.font.size = Pt(48)
    run.font.color.rgb = text_color
    run.font.bold = True
    set_run_complex_font(run, "Suez One")